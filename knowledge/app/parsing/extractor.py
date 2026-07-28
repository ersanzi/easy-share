"""多格式文档解析：TXT/Markdown/DOCX/文本型 PDF/XLSX/PPTX → 统一结构。"""
from __future__ import annotations

import io
import logging
import os
import re
import time
import zipfile
from collections.abc import Iterable

from app.parsing.cleaner import clean_document
from app.parsing.models import DocumentBlock, ParsedDocument, SourceLocation
from app.parsing.renderer import render_markdown
from app.ocr import (
    OCRPageResult,
    OCRProvider,
    OCRRecognitionError,
    OCRUnavailableError,
    UnavailableOCRProvider,
)

logger = logging.getLogger(__name__)
SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".docx", ".pdf", ".xlsx", ".pptx",
    ".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff",
}
OLE_COMPOUND_FILE_SIGNATURE = bytes.fromhex("D0CF11E0A1B11AE1")
OFFICE_OPEN_XML_FORMATS = {
    ".docx": (".doc", "word/document.xml"),
    ".xlsx": (".xls", "xl/workbook.xml"),
    ".pptx": (".ppt", "ppt/presentation.xml"),
}


class DocumentParseError(ValueError):
    """文件格式不支持、损坏或没有可提取内容。"""


def _validate_office_open_xml(filename: str, extension: str, content: bytes) -> None:
    """在第三方解析器前识别旧版 Office、损坏容器和 OOXML 类型错配。"""
    legacy_extension, required_member = OFFICE_OPEN_XML_FORMATS[extension]
    if content.startswith(OLE_COMPOUND_FILE_SIGNATURE):
        raise DocumentParseError(
            f"{filename} 的内容是旧版 Office 二进制格式（{legacy_extension}），"
            f"与扩展名 {extension} 不一致；请用 Word/WPS 打开后“另存为” {extension} 再上传"
        )

    try:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            members = set(archive.namelist())
    except zipfile.BadZipFile as exc:
        raise DocumentParseError(
            f"{filename} 不是有效的 Office Open XML 文件；"
            "可能已损坏、下载不完整或扩展名被修改"
        ) from exc

    if required_member in members:
        return

    for actual_extension, (_, core_member) in OFFICE_OPEN_XML_FORMATS.items():
        if core_member in members:
            raise DocumentParseError(
                f"{filename} 的实际内容是 {actual_extension}，与扩展名 {extension} 不一致；"
                f"请使用正确的文件名，或用 Office/WPS 另存为 {extension} 后再上传"
            )

    raise DocumentParseError(
        f"{filename} 不是有效的 {extension} 文件：缺少必要核心结构 {required_member}；"
        "文件可能已损坏或扩展名被修改"
    )


def parse_document(
    filename: str,
    content: bytes,
    *,
    ocr_provider: OCRProvider | None = None,
    ocr_min_text_chars: int = 20,
) -> ParsedDocument:
    extension = os.path.splitext(filename.lower())[1]
    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise DocumentParseError(f"不支持的文件格式 {extension or '(无扩展名)'}，当前支持: {supported}")
    if not content:
        raise DocumentParseError("文件内容为空")
    if extension in OFFICE_OPEN_XML_FORMATS:
        _validate_office_open_xml(filename, extension, content)

    try:
        if extension in {".txt", ".md", ".markdown"}:
            document = _parse_text(filename, content, markdown=extension != ".txt")
        elif extension == ".docx":
            document = _parse_docx(filename, content)
        elif extension == ".pdf":
            document = _parse_pdf(
                filename,
                content,
                ocr_provider=ocr_provider,
                ocr_min_text_chars=ocr_min_text_chars,
            )
        elif extension in {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}:
            document = _parse_image(filename, content, ocr_provider=ocr_provider)
        elif extension == ".xlsx":
            document = _parse_xlsx(filename, content)
        else:
            document = _parse_pptx(filename, content)
    except DocumentParseError:
        raise
    except Exception as exc:
        raise DocumentParseError(f"解析 {filename} 失败: {exc}") from exc

    document = clean_document(document)
    if not document.blocks:
        raise DocumentParseError(f"{filename} 没有可提取的文本或表格内容")
    return document


def extract_text(filename: str, content: bytes) -> str:
    """兼容旧入口：新解析器完成结构化解析后渲染为 Markdown 文本。"""
    return render_markdown(parse_document(filename, content)).strip()


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8-sig", "gb18030"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise DocumentParseError("文本文件编码无法识别，仅支持 UTF-8 或 GB18030 兼容编码")


def _parse_text(filename: str, content: bytes, *, markdown: bool) -> ParsedDocument:
    text = _decode_text(content)
    if markdown:
        return _parse_markdown(filename, text)

    blocks: list[DocumentBlock] = []
    for index, raw in enumerate(re.split(r"\n\s*\n", text), start=1):
        value = raw.strip()
        if not value:
            continue
        blocks.append(
            DocumentBlock(
                id=f"b{len(blocks) + 1}",
                type="paragraph",
                text=value,
                source=SourceLocation(paragraph=index),
            )
        )
    return ParsedDocument(filename=filename, media_type="text/plain", blocks=blocks)


def _parse_markdown(filename: str, text: str) -> ParsedDocument:
    """按标题边界切分 Markdown，避免标题与紧随正文被合并成普通段落。"""
    blocks: list[DocumentBlock] = []
    paragraph_lines: list[str] = []
    paragraph_start = 1

    def flush_paragraph() -> None:
        if not paragraph_lines:
            return
        blocks.append(
            DocumentBlock(
                id=f"b{len(blocks) + 1}",
                type="paragraph",
                text="\n".join(paragraph_lines).strip(),
                source=SourceLocation(paragraph=paragraph_start),
            )
        )
        paragraph_lines.clear()

    for line_number, raw_line in enumerate(text.splitlines(), start=1):
        heading = re.fullmatch(r"\s*(#{1,6})\s+(.+?)\s*", raw_line)
        if heading:
            flush_paragraph()
            blocks.append(
                DocumentBlock(
                    id=f"b{len(blocks) + 1}",
                    type="heading",
                    text=heading.group(2),
                    level=len(heading.group(1)),
                    source=SourceLocation(paragraph=line_number),
                )
            )
            continue
        if not raw_line.strip():
            flush_paragraph()
            continue
        if not paragraph_lines:
            paragraph_start = line_number
        paragraph_lines.append(raw_line)

    flush_paragraph()
    return ParsedDocument(filename=filename, media_type="text/markdown", blocks=blocks)

def _iter_docx_blocks(document) -> Iterable[tuple[str, object]]:
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P

    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield "paragraph", Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield "table", Table(child, document)


def _parse_docx(filename: str, content: bytes) -> ParsedDocument:
    from docx import Document

    document = Document(io.BytesIO(content))
    blocks: list[DocumentBlock] = []
    paragraph_number = 0
    table_number = 0
    for kind, item in _iter_docx_blocks(document):
        if kind == "paragraph":
            paragraph_number += 1
            text = item.text.strip()
            if not text:
                continue
            style_name = (item.style.name if item.style is not None else "") or ""
            match = re.search(r"(?:Heading|标题)\s*([1-6])", style_name, re.IGNORECASE)
            block_type = "heading" if match else "paragraph"
            blocks.append(
                DocumentBlock(
                    id=f"b{len(blocks) + 1}",
                    type=block_type,
                    text=text,
                    level=int(match.group(1)) if match else None,
                    source=SourceLocation(paragraph=paragraph_number),
                    metadata={"style": style_name} if style_name else {},
                )
            )
        else:
            table_number += 1
            rows = [[cell.text.strip() for cell in row.cells] for row in item.rows]
            blocks.append(
                DocumentBlock(
                    id=f"b{len(blocks) + 1}",
                    type="table",
                    rows=rows,
                    source=SourceLocation(table=table_number),
                )
            )
    return ParsedDocument(
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        blocks=blocks,
        metadata={"paragraphs": paragraph_number, "tables": table_number},
    )


def _render_pdf_page(content: bytes, page_number: int) -> bytes:
    """将 PDF 单页渲染为 PNG，供 OCR 识别。"""
    try:
        import fitz
    except ImportError as exc:
        raise DocumentParseError(
            "扫描 PDF OCR 需要 PyMuPDF；请安装 knowledge/requirements-ocr.txt"
        ) from exc
    try:
        document = fitz.open(stream=content, filetype="pdf")
        try:
            page = document.load_page(page_number - 1)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            return pixmap.tobytes("png")
        finally:
            document.close()
    except Exception as exc:  # noqa: BLE001 - 统一转换为可诊断的领域错误
        raise DocumentParseError(f"PDF 第 {page_number} 页渲染为 OCR 图片失败：{exc}") from exc


def _ocr_blocks(result: OCRPageResult, *, page: int) -> list[DocumentBlock]:
    blocks: list[DocumentBlock] = []
    for index, item in enumerate(result.blocks, start=1):
        text = item.text.strip()
        if not text:
            continue
        metadata = {"extraction_method": "ocr"}
        if item.confidence is not None:
            metadata["ocr_confidence"] = item.confidence
        if item.bbox:
            metadata["ocr_bbox"] = item.bbox
        blocks.append(
            DocumentBlock(
                id=f"ocr-p{page}-{index}",
                type="paragraph",
                text=text,
                source=SourceLocation(page=page),
                metadata=metadata,
            )
        )
    return blocks


def _resolve_ocr_provider(ocr_provider: OCRProvider | None) -> OCRProvider:
    return ocr_provider or UnavailableOCRProvider()


def _text_layer_blocks(paragraphs: list[str], *, page: int, offset: int) -> list[DocumentBlock]:
    return [
        DocumentBlock(
            id=f"b{offset + index}",
            type="paragraph",
            text=paragraph,
            source=SourceLocation(page=page),
        )
        for index, paragraph in enumerate(paragraphs, start=1)
    ]


def _parse_pdf(
    filename: str,
    content: bytes,
    *,
    ocr_provider: OCRProvider | None,
    ocr_min_text_chars: int,
) -> ParsedDocument:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    provider = _resolve_ocr_provider(ocr_provider)
    capability = provider.capability()
    blocks: list[DocumentBlock] = []
    empty_pages = 0
    ocr_pages: list[int] = []
    failed_pages: list[int] = []
    low_confidence_blocks = 0
    ocr_duration_ms = 0

    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        paragraphs = [
            part.strip()
            for part in re.split(r"\n\s*\n|(?<=。)\s*\n", text)
            if part.strip()
        ]
        text_blocks = _text_layer_blocks(paragraphs, page=page_number, offset=len(blocks))
        if not paragraphs:
            empty_pages += 1
        if paragraphs and len(text.strip()) >= ocr_min_text_chars:
            blocks.extend(text_blocks)
            continue

        if not capability.available:
            if text_blocks:
                blocks.extend(text_blocks)
            else:
                failed_pages.append(page_number)
            continue

        try:
            rendered = _render_pdf_page(content, page_number)
            result = provider.recognize_image(rendered, filename=filename, page=page_number)
        except OCRUnavailableError:
            if text_blocks:
                blocks.extend(text_blocks)
            else:
                failed_pages.append(page_number)
            continue
        except OCRRecognitionError:
            if text_blocks:
                blocks.extend(text_blocks)
            else:
                failed_pages.append(page_number)
            continue

        page_blocks = _ocr_blocks(result, page=page_number)
        if page_blocks:
            blocks.extend(page_blocks)
            ocr_pages.append(page_number)
            low_confidence_blocks += result.low_confidence_blocks
            ocr_duration_ms += result.duration_ms or 0
        elif text_blocks:
            blocks.extend(text_blocks)
        else:
            failed_pages.append(page_number)

    if not blocks:
        reason = capability.reason or "OCR 未识别出有效文字"
        raise DocumentParseError(
            f"{filename} 没有可提取的文本层；扫描件 OCR 不可用：{reason}。"
            "请安装并初始化 PaddleOCR，或上传带文本层的 PDF"
        )

    warnings: list[str] = []
    if failed_pages:
        warnings.append(
            f"有 {len(failed_pages)} 页未提取到文本，可能包含扫描页或纯图片；"
            f"OCR 未成功：{failed_pages}"
        )
    metadata = {
        "pages": len(reader.pages),
        "empty_pages": empty_pages,
        "ocr": {
            "provider": capability.provider,
            "pages_total": len(reader.pages),
            "pages_ocr": ocr_pages,
            "failed_pages": failed_pages,
            "low_confidence_blocks": low_confidence_blocks,
            "duration_ms": ocr_duration_ms,
        },
    }
    return ParsedDocument(
        filename=filename,
        media_type="application/pdf",
        blocks=blocks,
        warnings=warnings,
        metadata=metadata,
    )


def _parse_image(filename: str, content: bytes, *, ocr_provider: OCRProvider | None) -> ParsedDocument:
    provider = _resolve_ocr_provider(ocr_provider)
    started = time.perf_counter()
    try:
        result = provider.recognize_image(content, filename=filename, page=1)
    except OCRUnavailableError as exc:
        raise DocumentParseError(
            f"{filename} 需要 OCR 才能解析，但当前不可用：{exc}。"
            "请安装 knowledge/requirements-ocr.txt 后重试"
        ) from exc
    except OCRRecognitionError as exc:
        raise DocumentParseError(str(exc)) from exc
    blocks = _ocr_blocks(result, page=1)
    if not blocks:
        raise DocumentParseError(f"{filename} OCR 未识别出有效文字，请检查图片清晰度")
    duration_ms = result.duration_ms or round((time.perf_counter() - started) * 1000)
    return ParsedDocument(
        filename=filename,
        media_type={
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".bmp": "image/bmp",
            ".tif": "image/tiff",
            ".tiff": "image/tiff",
        }[os.path.splitext(filename.lower())[1]],
        blocks=blocks,
        metadata={
            "pages": 1,
            "ocr": {
                "provider": provider.capability().provider,
                "pages_total": 1,
                "pages_ocr": [1],
                "failed_pages": [],
                "low_confidence_blocks": result.low_confidence_blocks,
                "duration_ms": duration_ms,
            },
        },
    )


def _cell_text(value: object) -> str:
    if value is None:
        return ""
    if hasattr(value, "isoformat"):
        try:
            return value.isoformat()
        except TypeError:
            pass
    return str(value)


def _parse_xlsx(filename: str, content: bytes) -> ParsedDocument:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    blocks: list[DocumentBlock] = []
    sheet_count = len(workbook.sheetnames)
    try:
        for worksheet in workbook.worksheets:
            blocks.append(
                DocumentBlock(
                    id=f"b{len(blocks) + 1}",
                    type="heading",
                    text=worksheet.title,
                    level=1,
                    source=SourceLocation(sheet=worksheet.title),
                )
            )
            rows: list[list[str]] = []
            start_row: int | None = None
            for row_number, values in enumerate(worksheet.iter_rows(values_only=True), start=1):
                cells = [_cell_text(value) for value in values]
                while cells and not cells[-1]:
                    cells.pop()
                if not any(cells):
                    if rows:
                        blocks.append(
                            DocumentBlock(
                                id=f"b{len(blocks) + 1}",
                                type="table",
                                rows=rows,
                                source=SourceLocation(sheet=worksheet.title, row=start_row),
                            )
                        )
                        rows = []
                        start_row = None
                    continue
                if start_row is None:
                    start_row = row_number
                rows.append(cells)
                if len(rows) >= 200:
                    blocks.append(
                        DocumentBlock(
                            id=f"b{len(blocks) + 1}",
                            type="table",
                            rows=rows,
                            source=SourceLocation(sheet=worksheet.title, row=start_row),
                        )
                    )
                    rows = []
                    start_row = None
            if rows:
                blocks.append(
                    DocumentBlock(
                        id=f"b{len(blocks) + 1}",
                        type="table",
                        rows=rows,
                        source=SourceLocation(sheet=worksheet.title, row=start_row),
                    )
                )
    finally:
        workbook.close()
    return ParsedDocument(
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        blocks=blocks,
        metadata={"sheets": sheet_count},
    )


def _parse_pptx(filename: str, content: bytes) -> ParsedDocument:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    blocks: list[DocumentBlock] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title_shape_id = title_shape.shape_id if title_shape is not None else None
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                blocks.append(
                    DocumentBlock(
                        id=f"b{len(blocks) + 1}",
                        type="table",
                        rows=rows,
                        source=SourceLocation(slide=slide_number),
                    )
                )
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph_number, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                text = paragraph.text.strip()
                if not text:
                    continue
                is_title = title_shape_id is not None and shape.shape_id == title_shape_id
                blocks.append(
                    DocumentBlock(
                        id=f"b{len(blocks) + 1}",
                        type="heading" if is_title else "paragraph",
                        text=text,
                        level=1 if is_title else None,
                        source=SourceLocation(slide=slide_number, paragraph=paragraph_number),
                    )
                )
    return ParsedDocument(
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        blocks=blocks,
        metadata={"slides": len(presentation.slides)},
    )
