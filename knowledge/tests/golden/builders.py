"""确定性构造 Office/PDF 黄金样本，避免提交来源不明的二进制文件。"""
from __future__ import annotations

import io
from collections.abc import Callable


def build_docx_policy() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("EasyShare 文档处理规范", level=1)
    document.add_paragraph("Office 文件必须保留结构和来源位置。")
    table = document.add_table(rows=3, cols=2)
    values = [
        ("字段", "含义"),
        ("file_id", "稳定文件身份"),
        ("version_id", "内容版本"),
    ]
    for row, values_row in zip(table.rows, values):
        for cell, value in zip(row.cells, values_row):
            cell.text = value
    document.add_heading("验收标准", level=2)
    document.add_paragraph("解析结果必须可追溯。")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def build_xlsx_project_register() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    projects = workbook.active
    projects.title = "项目清单"
    projects.append(["项目", "负责人", "状态"])
    projects.append(["文档清洗", "Alice", "完成"])
    projects.append(["扫描件 OCR", "Bob", "待开始"])

    risks = workbook.create_sheet("风险")
    risks.append(["风险", "等级"])
    risks.append(["表格跨页", "中"])

    buffer = io.BytesIO()
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def build_pptx_quality_review() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    overview = presentation.slides.add_slide(presentation.slide_layouts[1])
    overview.shapes.title.text = "季度知识库验收"
    overview.placeholders[1].text = "Office 文档结构解析完成"

    metrics = presentation.slides.add_slide(presentation.slide_layouts[5])
    metrics.shapes.title.text = "质量指标"
    table = metrics.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(8), Inches(2)).table
    values = [
        ("指标", "目标"),
        ("来源定位", "100%"),
        ("派生产物", "3 类"),
    ]
    for row_index, values_row in enumerate(values):
        for column_index, value in enumerate(values_row):
            table.cell(row_index, column_index).text = value

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def build_pdf_text_with_blank_page() -> bytes:
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    text_page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    text_page[NameObject("/Resources")] = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject(
                {NameObject("/F1"): font_ref}
            )
        }
    )
    stream = DecodedStreamObject()
    stream.set_data(b"BT /F1 18 Tf 72 720 Td (EasyShare PDF Golden Page One) Tj ET")
    text_page[NameObject("/Contents")] = writer._add_object(stream)
    writer.add_blank_page(width=612, height=792)

    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


BUILDERS: dict[str, Callable[[], bytes]] = {
    "docx_policy": build_docx_policy,
    "xlsx_project_register": build_xlsx_project_register,
    "pptx_quality_review": build_pptx_quality_review,
    "pdf_text_with_blank_page": build_pdf_text_with_blank_page,
}


def build_case(builder_name: str) -> bytes:
    try:
        builder = BUILDERS[builder_name]
    except KeyError as exc:
        raise ValueError(f"未知黄金样本构造器: {builder_name}") from exc
    return builder()