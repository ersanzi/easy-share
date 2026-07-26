from __future__ import annotations

import io
import zipfile

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

from app.parsing.extractor import DocumentParseError, parse_document
from app.parsing.renderer import render_markdown


def test_txt_and_markdown_are_cleaned() -> None:
    txt = parse_document("note.txt", "first  \r\n\r\nfirst\r\n\r\nsecond\tline".encode())
    assert [block.text for block in txt.blocks] == ["first", "second line"]
    assert txt.warnings

    markdown = parse_document("guide.md", "# Title\nBody content\n## Details\nMore".encode())
    assert markdown.blocks[0].type == "heading"
    assert markdown.blocks[0].level == 1
    assert [block.type for block in markdown.blocks] == ["heading", "paragraph", "heading", "paragraph"]
    assert render_markdown(markdown).startswith("# Title")


def test_docx_keeps_heading_paragraph_and_table_order() -> None:
    source = Document()
    source.add_heading("Policy", level=1)
    source.add_paragraph("Body")
    table = source.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Team"
    table.cell(1, 0).text = "Easy"
    table.cell(1, 1).text = "R&D"
    buffer = io.BytesIO()
    source.save(buffer)

    parsed = parse_document("policy.docx", buffer.getvalue())
    assert [block.type for block in parsed.blocks] == ["heading", "paragraph", "table"]
    assert parsed.blocks[2].rows[1] == ["Easy", "R&D"]


def test_xlsx_preserves_sheet_and_row_location() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Employees"
    sheet.append(["Name", "Team"])
    sheet.append(["Easy", "R&D"])
    buffer = io.BytesIO()
    workbook.save(buffer)

    parsed = parse_document("employees.xlsx", buffer.getvalue())
    assert parsed.metadata["sheets"] == 1
    assert parsed.blocks[0].text == "Employees"
    assert parsed.blocks[1].source.sheet == "Employees"
    assert parsed.blocks[1].source.row == 1
    assert parsed.blocks[1].rows[1] == ["Easy", "R&D"]


def test_pptx_preserves_slide_content() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "Document cleaning completed"
    buffer = io.BytesIO()
    presentation.save(buffer)

    parsed = parse_document("review.pptx", buffer.getvalue())
    assert parsed.metadata["slides"] == 1
    assert any(block.text == "Quarterly Review" and block.source.slide == 1 for block in parsed.blocks)
    assert any("Document cleaning completed" in block.text for block in parsed.blocks)


def test_blank_pdf_requires_ocr() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    buffer = io.BytesIO()
    writer.write(buffer)

    with pytest.raises(DocumentParseError, match="OCR"):
        parse_document("scan.pdf", buffer.getvalue())


@pytest.mark.parametrize(
    ("filename", "legacy_extension", "modern_extension"),
    [
        ("员工手册.docx", ".doc", ".docx"),
        ("预算.xlsx", ".xls", ".xlsx"),
        ("汇报.pptx", ".ppt", ".pptx"),
    ],
)
def test_legacy_office_with_modern_extension_has_actionable_error(
    filename: str,
    legacy_extension: str,
    modern_extension: str,
) -> None:
    legacy_content = bytes.fromhex("D0CF11E0A1B11AE1") + b"\x00" * 64

    with pytest.raises(DocumentParseError) as exc_info:
        parse_document(filename, legacy_content)

    message = str(exc_info.value)
    assert "旧版 Office 二进制格式" in message
    assert legacy_extension in message
    assert modern_extension in message
    assert "另存为" in message
    assert "File is not a zip file" not in message


@pytest.mark.parametrize("filename", ["bad.docx", "bad.xlsx", "bad.pptx"])
def test_corrupted_office_file_fails_explicitly(filename: str) -> None:
    with pytest.raises(DocumentParseError) as exc_info:
        parse_document(filename, b"not-an-office-file")

    message = str(exc_info.value)
    assert filename in message
    assert "不是有效的 Office Open XML 文件" in message
    assert "损坏" in message
    assert "File is not a zip file" not in message


def test_office_open_xml_extension_mismatch_is_reported() -> None:
    workbook = Workbook()
    workbook.active.append(["Name", "Team"])
    buffer = io.BytesIO()
    workbook.save(buffer)

    with pytest.raises(DocumentParseError) as exc_info:
        parse_document("employees.docx", buffer.getvalue())

    message = str(exc_info.value)
    assert "实际内容是 .xlsx" in message
    assert "扩展名 .docx 不一致" in message


def test_office_open_xml_missing_core_member_is_reported() -> None:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types />")

    with pytest.raises(DocumentParseError) as exc_info:
        parse_document("empty.docx", buffer.getvalue())

    message = str(exc_info.value)
    assert "缺少必要核心结构 word/document.xml" in message


def test_unsupported_format_fails() -> None:
    with pytest.raises(DocumentParseError):
        parse_document("archive.zip", b"data")
